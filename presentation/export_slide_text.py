#!/usr/bin/env python3
"""Dump every slide's text as a PowerPoint outline you can paste in.

PowerPoint's Outline view imports tab-indented plain text: a line at column
zero starts a slide and becomes its title, each tab indents one level. So this
is the shortest path from the generated deck into someone else's template -
paste it into Outline view and the structure arrives, then apply the template.

Reading order is (top, then left), which is how the slides are laid out, and
the title is the largest run on the slide rather than the topmost - the eyebrow
sits above the title but is not it.

    python3 presentation/export_slide_text.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

HERE = Path(__file__).resolve().parent
DECK = HERE / "AI-Storytelling-Final-Talk.pptx"
OUT = HERE / "SLIDE-TEXT.txt"
EMU_IN = 914400.0


def biggest(shape) -> float:
    return max((r.font.size.pt for p in shape.text_frame.paragraphs
                for r in p.runs if r.font.size), default=0.0)


def main() -> int:
    prs = Presentation(DECK)
    lines: list[str] = []
    for n, slide in enumerate(prs.slides, 1):
        blocks = [sh for sh in slide.shapes
                  if sh.has_text_frame and sh.text_frame.text.strip()]
        if not blocks:                       # a full-bleed figure
            fig = next((sh for sh in slide.shapes if sh.shape_type == 13), None)
            lines.append(f"[Slide {n}] Figure: "
                         f"{Path(fig.image.filename).name if fig and fig.image.filename else 'image'}")
            note = slide.notes_slide.notes_text_frame.text.strip().splitlines()
            if note:
                lines.append(f"\t{note[0]}")
            lines.append("")
            continue

        title_shape = max(blocks, key=biggest)
        rest = sorted((b for b in blocks if b is not title_shape),
                      key=lambda s: (round(s.top / EMU_IN, 1), s.left))

        title = next((ln for ln in title_shape.text_frame.text.splitlines()
                      if ln.strip()), f"Slide {n}")
        lines.append(title.strip())
        for para in title_shape.text_frame.text.splitlines()[1:]:
            if para.strip():
                lines.append(f"\t{para.strip()}")
        for b in rest:
            body = [ln.strip() for ln in b.text_frame.text.splitlines() if ln.strip()]
            if not body:
                continue
            lines.append(f"\t{body[0]}")
            for ln in body[1:]:
                lines.append(f"\t\t{ln}")
        lines.append("")

    OUT.write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(f"wrote {OUT.name}  ({len(prs.slides)} slides, "
          f"{len([l for l in lines if l])} lines)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
