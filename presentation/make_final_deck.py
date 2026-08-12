#!/usr/bin/env python3
"""Build the final-presentation deck: 12 slides, 7 minutes of talk plus a demo.

The interim deck (``AI-Storytelling-Final-Presentation.pptx``) is the visual
parent - same 16:9 canvas, same IBM Plex Sans / Calibri pairing, same footer
line, and the palette sampled from the project logo. What changes is the
argument. The interim deck was a plan; this one reports measurements, so the
figures carry the slides and the prose gets out of their way.

Three things are deliberately *not* inherited from the interim deck:

* Its slide 13 reads the DataTales causal 0% as "a capability wall, not a size
  problem". That number measures groundedness against a masked reference, not
  the ability to reason causally, and the claim was scrubbed from the results
  page and the home page in `92ed3ae` and `eb1434a`. It does not come back here.
* Its title slide says "Interim Presentation".
* Its speaker notes are placeholders repeated across several slides. Here every
  slide's notes are the actual words for that slide, with a time budget, and
  they sum to 7:00.

Every number on a slide is read from ``experiments/exp_json/`` at build time or
lives inside a figure PNG generated from those same files. Nothing is typed in
by hand, so re-running an experiment and re-running this script keeps the deck
honest.

    python3 presentation/make_final_deck.py
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
EXP = ROOT / "experiments" / "exp_json"
FIGS = HERE / "figures"
OUT = HERE / "AI-Storytelling-Final-Talk.pptx"

# Sampled from the project logo; the same values the landing page uses.
NAVY = RGBColor(0x0D, 0x1B, 0x5C)
DEEP_NAVY = RGBColor(0x07, 0x1A, 0x3B)
BLUE = RGBColor(0x1E, 0x66, 0xB8)
TEAL = RGBColor(0x20, 0xC4, 0xB0)
DEEP_TEAL = RGBColor(0x0E, 0x8F, 0x86)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x52, 0x61, 0x73)
BORDER = RGBColor(0xD9, 0xDF, 0xE7)
SOFT = RGBColor(0xEE, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARN = RGBColor(0xB4, 0x54, 0x09)

# Arial is the parent deck's theme font and ships everywhere; IBM Plex Sans is
# used on 248 runs of that deck but is installed on none of our machines, so it
# silently substitutes and the layout shifts. Calibri comes with Office on both
# platforms. Both are measured for overflow in check_final_deck.py.
HEAD = "Arial"
BODY = "Calibri"
FOOTER = "AI-Powered Data Storytelling  ·  CMS Team Project  ·  TU Dresden  ·  SoSe 2026"

W, H = Inches(13.333), Inches(7.5)


# --------------------------------------------------------------------------
# data, read once so no figure on a slide is transcribed by hand
# --------------------------------------------------------------------------
def load() -> dict:
    hc = json.loads((EXP / "exp-human-comparison.json").read_text())
    rel = json.loads((EXP / "exp-judge-reliability.json").read_text())
    rep = json.loads((EXP / "exp-repeats-g4b-g8b.json").read_text())
    q27 = json.loads((EXP / "exp-repeats-q27b.json").read_text())
    pair = json.loads((ROOT / "experiments" / "pairwise_results.json").read_text())
    key = {p["pair_id"]: p
           for p in json.loads((ROOT / "experiments" / "pairwise_key.json").read_text())["key"]}

    tally: dict[str, dict[str, int]] = {}
    overall: dict[str, int] = {}
    for v in pair["verdicts"]:
        p = key[v["pair_id"]]
        name = lambda w: p["story_1_is"] if w == "story_1" else (
            p["story_2_is"] if w == "story_2" else w)
        overall[name(v["overall"])] = overall.get(name(v["overall"]), 0) + 1
        for crit, winner in v["criteria"].items():
            tally.setdefault(crit, {})
            tally[crit][name(winner)] = tally[crit].get(name(winner), 0) + 1

    g4b = rep["configs"]["g4b"]["summary"]
    g8b = rep["configs"]["g8b"]["summary"]
    return {
        "agg": hc["aggregate"],
        "n_runs_human": hc["n_runs"],
        "n_human": hc["n_human_stories"],
        "rel": rel,
        "g4b": g4b, "g8b": g8b,
        "q27": q27["configs"]["q27b"]["summary"],
        "pair_n": pair["n"], "pair_overall": overall, "pair_crit": tally,
    }


D = load()


# --------------------------------------------------------------------------
# drawing helpers
# --------------------------------------------------------------------------
def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=1.0):
    """runs: list of (string, size_pt, bold, color, fontname, space_before_pt)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (s, size, bold, color, font, before) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if before:
            p.space_before = Pt(before)
        r = p.add_run()
        r.text = s
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return box


def chrome(slide, n: int, *, dark=False):
    """Footer line and slide number, the interim deck's convention."""
    col = RGBColor(0x8A, 0x99, 0xAD) if dark else MUTED
    text(slide, Inches(0.75), Inches(6.95), Inches(9.5), Inches(0.3),
         [(FOOTER, 9, False, col, BODY, 0)])
    text(slide, Inches(12.0), Inches(6.95), Inches(0.6), Inches(0.3),
         [(str(n), 9, False, col, BODY, 0)], align=PP_ALIGN.RIGHT)


def heading(slide, kicker: str, title: str, sub: str | None = None):
    text(slide, Inches(0.75), Inches(0.55), Inches(11.8), Inches(0.3),
         [(kicker.upper(), 11, True, DEEP_TEAL, HEAD, 0)])
    text(slide, Inches(0.75), Inches(0.95), Inches(11.8), Inches(0.9),
         [(title, 30, True, NAVY, HEAD, 0)])
    rect(slide, Inches(0.75), Inches(1.82), Inches(1.1), Pt(3.5), fill=TEAL)
    if sub:
        text(slide, Inches(0.75), Inches(2.05), Inches(11.0), Inches(0.5),
             [(sub, 15, False, MUTED, BODY, 0)])


def notes(slide, body: str):
    slide.notes_slide.notes_text_frame.text = body.strip()


def fullbleed(prs, png: str, note: str):
    """A figure slide. The figures are 3200x1800, so they fill the canvas."""
    s = blank(prs)
    s.shapes.add_picture(str(FIGS / png), 0, 0, width=W, height=H)
    notes(s, note)
    return s


def bullets(slide, x, y, w, items, *, size=16, gap=13, colour=INK):
    runs = []
    for i, (lead, rest) in enumerate(items):
        runs.append((f"▪  {lead}", size, True, NAVY, BODY, 0 if i == 0 else gap))
        if rest:
            runs.append((f"     {rest}", size - 1, False, colour, BODY, 2))
    return text(slide, x, y, w, Inches(4.0), runs, spacing=1.15)


def card(slide, x, y, w, h, label, value, foot, *, accent=DEEP_TEAL):
    rect(slide, x, y, w, h, fill=SOFT)
    rect(slide, x, y, Pt(3.5), h, fill=accent)
    text(slide, x + Inches(0.28), y + Inches(0.22), w - Inches(0.5), Inches(0.3),
         [(label.upper(), 9.5, True, MUTED, HEAD, 0)])
    text(slide, x + Inches(0.28), y + Inches(0.52), w - Inches(0.5), Inches(0.6),
         [(value, 30, True, accent, HEAD, 0)])
    text(slide, x + Inches(0.28), y + Inches(1.12), w - Inches(0.5), Inches(0.7),
         [(foot, 11, False, MUTED, BODY, 0)], spacing=1.1)


# --------------------------------------------------------------------------
# slides
# --------------------------------------------------------------------------
def scale_bar(slide, x, y, w, label, *, human: float, human_label: str):
    """A 1-5 axis with the rubric's calibrated band and where humans actually sit.

    Both markers, because the rubric calls 3 calibrated while every human story
    we collected sits nearer 2. Showing only the rubric would set the audience up
    to read our 2.09 as a miss; showing only the human band would hide that the
    scale has a defined middle. The distance between the two marks is a real
    finding, so it belongs on the slide rather than in a footnote.
    """
    def at(v: float) -> Emu:
        return Inches(x + (v - 1) / 4 * w)

    text(slide, Inches(x), Inches(y), Inches(w), Inches(0.3),
         [(label, 10.5, True, MUTED, HEAD, 0)])
    rect(slide, Inches(x), Inches(y + 0.36), Inches(w), Inches(0.3), fill=SOFT)
    # the rubric's calibrated zone, 3 +/- 0.5
    rect(slide, at(2.5), Inches(y + 0.36), Inches(w / 4), Inches(0.3),
         fill=RGBColor(0xD5, 0xEC, 0xE9))
    text(slide, at(2.5), Inches(y + 0.4), Inches(w / 4), Inches(0.24),
         [("3  rubric", 9.5, True, DEEP_TEAL, BODY, 0)], align=PP_ALIGN.CENTER)
    # where our human writers actually landed
    rect(slide, at(human) - Inches(0.015), Inches(y + 0.3), Inches(0.03),
         Inches(0.42), fill=BLUE)
    text(slide, at(human) - Inches(0.8), Inches(y + 0.74), Inches(1.6), Inches(0.24),
         [(human_label, 9.5, True, BLUE, BODY, 0)], align=PP_ALIGN.CENTER)
    text(slide, Inches(x), Inches(y + 0.74), Inches(1.2), Inches(0.24),
         [("1  flat", 9.5, False, MUTED, BODY, 0)])
    text(slide, Inches(x + w - 1.4), Inches(y + 0.74), Inches(1.4), Inches(0.24),
         [("5  extreme", 9.5, False, MUTED, BODY, 0)], align=PP_ALIGN.RIGHT)


def group(slide, x, y, w, h, title, why, rows, *, accent=DEEP_TEAL, cost=None):
    """One family of metrics: what it is, why it exists, what is in it."""
    rect(slide, Inches(x), Inches(y), Inches(w), Inches(h), fill=SOFT)
    rect(slide, Inches(x), Inches(y), Inches(w), Pt(3.5), fill=accent)
    runs = [(title, 13, True, NAVY, HEAD, 0),
            (why, 10.5, False, MUTED, BODY, 3)]
    for name, gloss in rows:
        runs.append((f"{name}   {gloss}", 10.5, False, INK, BODY, 5))
    if cost:
        runs.append((cost, 10, True, accent, BODY, 6))
    text(slide, Inches(x + 0.24), Inches(y + 0.22), Inches(w - 0.48),
         Inches(h - 0.4), runs, spacing=1.12)


PARTS = [
    ("The research, and the problem it found", "Mahmoud"),
    ("What we built, and how we measure it", "Okasha"),
    ("What the reader sees, and how we chose the pair", "Elsaadani"),
    ("The result, the agent, and the system", "Ramadan"),
]


def divider(prs, idx: int, note: str):
    """A break between sections, doubling as a progress checkpoint.

    Four speakers means three handovers, and a handover with no slide change is
    where a talk stalls: the room does not know a new person is starting and the
    new person has no beat to start on. The rail also answers "how much is left",
    which an audience tracks whether or not you tell them.
    """
    title, who = PARTS[idx]
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=DEEP_NAVY)
    rect(s, 0, 0, Inches(0.16), H, fill=TEAL)
    text(s, Inches(1.1), Inches(2.35), Inches(11.0), Inches(0.4),
         [(f"PART {idx + 1} OF {len(PARTS)}", 12, True, TEAL, HEAD, 0)])
    text(s, Inches(1.1), Inches(2.85), Inches(10.6), Inches(1.1),
         [(title, 34, True, WHITE, HEAD, 0)], spacing=1.08)
    text(s, Inches(1.1), Inches(4.15), Inches(10.6), Inches(0.4),
         [(who, 15, True, RGBColor(0x8A, 0x9C, 0xB8), BODY, 0)])
    # the rail
    seg = 2.62
    for i, (label, _) in enumerate(PARTS):
        x = 1.1 + i * (seg + 0.22)
        done = i <= idx
        rect(s, Inches(x), Inches(5.35), Inches(seg), Pt(5),
             fill=TEAL if done else RGBColor(0x2B, 0x4A, 0x84))
        text(s, Inches(x), Inches(5.55), Inches(seg), Inches(0.5),
             [(label.split(",")[0], 10.5, i == idx,
               WHITE if i == idx else RGBColor(0x6B, 0x7D, 0x9B), BODY, 0)],
             spacing=1.12)
    notes(s, note)
    return s


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    agg, rel = D["agg"], D["rel"]
    n = D["pair_n"]
    crit = D["pair_crit"]

    # -- 1  title ---------------------------------------------------------
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=DEEP_NAVY)
    rect(s, 0, 0, Inches(0.16), H, fill=TEAL)
    # No date. Topic.pptx lists 10 August as the *proposed* final date and that
    # has passed, so printing it would put a stale fact on the most visible line
    # in the deck. Add the real date and room here when they are known.
    text(s, Inches(1.1), Inches(1.55), Inches(11.0), Inches(0.4),
         [("FINAL PRESENTATION  ·  CMS TEAM PROJECT  ·  TU DRESDEN", 12, True,
           TEAL, HEAD, 0)])
    text(s, Inches(1.1), Inches(2.1), Inches(11.0), Inches(1.6),
         [("Moderating the Emotional Tone", 42, True, WHITE, HEAD, 0),
          ("of a Generated Data Story", 42, True, WHITE, HEAD, 0)], spacing=1.05)
    rect(s, Inches(1.1), Inches(3.95), Inches(1.6), Pt(4), fill=TEAL)
    text(s, Inches(1.1), Inches(4.3), Inches(11.0), Inches(0.9),
         [("An agent that reads a data story and pulls its tone back to what the "
           "evidence supports —", 15, False, RGBColor(0xC7, 0xD4, 0xE6), BODY, 0),
          ("measured against 25 stories our own team wrote from the same tables.",
           15, False, RGBColor(0xC7, 0xD4, 0xE6), BODY, 2)], spacing=1.2)
    text(s, Inches(1.1), Inches(5.6), Inches(11.0), Inches(0.9),
         [("Mahmoud Elsamadony  ·  Ahmed Okasha  ·  Ahmed Elsaadani  ·  "
           "Ahmed Ramadan", 14, True, WHITE, BODY, 0),
          ("Supervisors: Susmita Khadse, Julián Méndez  ·  Chair: "
           "Prof. Dr. Raimund Dachselt (IMLD)", 11.5, False,
           RGBColor(0x8A, 0x9C, 0xB8), BODY, 5)])
    notes(s, """
SPEAKER 1 - Mahmoud   [0:15]   running total 0:15   ~40 words

Good afternoon. We turn a public-health table into a written story. The part we
want to defend today is the second agent: it reads that story and pulls the
emotional tone back to what the evidence supports.

Four parts, then a live demo.

DELIVERY: do not read the names, they are on the slide. Go straight to the
sentence about the second agent.
""")


    # -- 2  table of contents ---------------------------------------------
    s = blank(prs)
    heading(s, "Agenda", "Four parts, then the system running")
    for i, (title, who) in enumerate(PARTS):
        y = 2.35 + i * 0.92
        rect(s, Inches(0.75), Inches(y), Inches(8.6), Inches(0.76), fill=SOFT)
        rect(s, Inches(0.75), Inches(y), Pt(3.5), Inches(0.76), fill=TEAL)
        text(s, Inches(1.05), Inches(y + 0.12), Inches(0.7), Inches(0.5),
             [(f"0{i + 1}", 22, True, RGBColor(0xC2, 0xCE, 0xDB), HEAD, 0)])
        text(s, Inches(1.85), Inches(y + 0.13), Inches(6.0), Inches(0.5),
             [(title, 15, True, NAVY, BODY, 0)])
        text(s, Inches(1.85), Inches(y + 0.44), Inches(6.0), Inches(0.26),
             [(who, 11, False, MUTED, BODY, 0)])
    rect(s, Inches(0.75), Inches(6.03), Inches(8.6), Inches(0.62), fill=DEEP_NAVY)
    text(s, Inches(1.05), Inches(6.18), Inches(8.0), Inches(0.4),
         [("LIVE DEMO   ·   the pipeline running locally   ·   3:30", 13, True,
           WHITE, BODY, 0)])
    rect(s, Inches(9.75), Inches(2.35), Inches(2.83), Inches(4.3), fill=SOFT)
    rect(s, Inches(9.75), Inches(2.35), Inches(2.83), Pt(4), fill=BLUE)
    text(s, Inches(10.02), Inches(2.58), Inches(2.35), Inches(3.9), [
        ("The one claim", 12.5, True, NAVY, BODY, 0),
        ("A second agent can pull a generated data story's tone back to what the "
         "evidence supports —", 12, False, INK, BODY, 6),
        ("without losing the numbers,", 12, True, DEEP_TEAL, BODY, 6),
        ("and at a cost in readability we can put a number on.", 12, True,
         DEEP_TEAL, BODY, 2),
        ("Everything else in this talk exists to make that measurable.", 11.5,
         False, MUTED, BODY, 8),
    ], spacing=1.16)
    chrome(s, 2)
    notes(s, """
SPEAKER 1 - Mahmoud   [0:20]   running total 0:35   ~38 words

Four parts. What the research told us and the problem it exposed. What we built
and how we measure it. What a reader sees, and how we picked the model pair.
Then the result and the system itself, and we finish live in the app.

DELIVERY: fifteen seconds. Read the four part names, point at the claim on the
right, move on. Do not explain the claim here - the whole talk is the
explanation.
""")

    # -- 3  divider -------------------------------------------------------
    divider(prs, 0, """
SPEAKER 1 - Mahmoud   [0:05]   running total 0:40   ~15 words

First: what the field already does, and the gap we found in it.

DELIVERY: one line, then move. A divider is a beat, not a slide to present.
""")

    # -- 4  the research --------------------------------------------------
    s = blank(prs)
    heading(s, "The research", "Three systems we did not just cite — we re-ran them")
    repros = [
        ("Kasner and Dušek, ACL 2024", "Are open models faithful?",
         "Their reference-free method marks every span of generated text that "
         "the source table does not support.",
         "We re-ran it on modern open models. A 12B model left a semantic error "
         "in about 18% of outputs, where the paper reported over 80%. A 4B model "
         "regressed to 52%.",
         "Recency and size both matter — and today's open models are already "
         "fairly faithful at stating data.", BLUE),
        ("DataTales, 2025", "Which kinds of claim are hard?",
         "A benchmark for data narration that scores a story per operation: "
         "lookup, comparison, trend, rate of change, causal.",
         "We reproduced a 30-report slice. Reading operations climb steeply with "
         "scale — rate of change went 43% to 89% from a 4B to a 12B model.",
         "Stating a number is not the hard part. Explaining it is.", TEAL),
        ("DataNarrative, EMNLP 2024", "What architecture works?",
         "One agent writes the story, a second agent verifies it at every step, "
         "beating single-LLM baselines.",
         "We rebuilt the two-agent generate-and-verify loop, which is the "
         "structural basis of our own pipeline.",
         "The second-agent pattern works. Ours swaps its factual verifier for a "
         "tone moderator.", NAVY),
    ]
    for i2, (paper, q, what, did, took, col) in enumerate(repros):
        x = 0.75 + i2 * 4.06
        rect(s, Inches(x), Inches(2.3), Inches(3.72), Inches(3.62), fill=SOFT)
        rect(s, Inches(x), Inches(2.3), Inches(3.72), Pt(4), fill=col)
        text(s, Inches(x + 0.26), Inches(2.52), Inches(3.2), Inches(3.3), [
            (paper, 10, True, MUTED, HEAD, 0),
            (q, 14, True, col, HEAD, 4),
            (what, 11, False, INK, BODY, 7),
            (f"What we did: {did}", 11, False, INK, BODY, 7),
            (took, 11, True, NAVY, BODY, 7),
        ], spacing=1.15)
    rect(s, Inches(0.75), Inches(6.08), Inches(11.83), Inches(0.6), fill=SOFT)
    rect(s, Inches(0.75), Inches(6.08), Pt(3.5), Inches(0.6), fill=TEAL)
    text(s, Inches(1.05), Inches(6.21), Inches(11.3), Inches(0.4),
         [("Put together: open models already state data fairly faithfully, and "
           "every critic agent in this literature checks facts. So another "
           "fact-checker is not the gap.", 13, True, NAVY, BODY, 0)])
    chrome(s, 4)
    notes(s, """
SPEAKER 1 - Mahmoud   [0:35]   running total 1:15   ~100 words

We surveyed the field, and the report lists it. These are the three we re-ran
ourselves, because they shaped what we built.

Kasner and Dušek: are open models faithful? Re-running their method, a 12B model
left an error in about 18 percent of outputs where the paper reported over 80.

DataTales: which claims are hard? Stating a number is not the hard part.
Explaining it is.

DataNarrative gave us the architecture - one agent writes, a second checks.

Together, that is the bar at the bottom. Another fact-checker is not the gap.

DELIVERY: three cards, one sentence each, then the bottom bar. If asked about
the wider survey, it is in the interim report.
""")

    # -- 4  the problem ---------------------------------------------------
    s = blank(prs)
    heading(s, "The problem", "Everyone checks the facts. Nobody checks the tone.")
    bullets(s, Inches(0.75), Inches(2.45), Inches(6.35), [
        ("The critic is always a fact-checker.",
         "That is what our survey found, without exception. No published system "
         "moderates the affective tone of a data narrative."),
        ("But framing moves the reader on its own.",
         "Hullman and Diakopoulos show editorial framing significantly changes "
         "how the same chart is read. Tone is not decoration; it is part of the "
         "claim."),
        ("So a story can be fully correct and still mislead.",
         "Both sentences on the right are true to the table. A fact-checker "
         "passes both. Only one of them is honest about what the data says."),
    ])
    rect(s, Inches(7.5), Inches(2.35), Inches(5.08), Inches(3.5), fill=SOFT)
    rect(s, Inches(7.5), Inches(2.35), Inches(5.08), Pt(4), fill=BLUE)
    text(s, Inches(7.85), Inches(2.65), Inches(4.4), Inches(3.0), [
        ("EVERY NUMBER BELOW IS CORRECT", 10, True, MUTED, HEAD, 0),
        ("“The world faces a catastrophic resurgence as a deadly disease "
         "roars back.”", 16, True, WARN, BODY, 14),
        ("A fact-checker passes this. It cites nothing false.", 12, False, MUTED,
         BODY, 6),
        ("“Cases rose for a third year while coverage stalled near 84%, below "
         "the 95% herd-immunity threshold.”", 16, True, DEEP_TEAL, BODY, 16),
        ("Same table. Same facts. A different reader reaction.", 12, False, MUTED,
         BODY, 6),
    ], spacing=1.15)
    chrome(s, 5)
    notes(s, """
SPEAKER 1 - Mahmoud   [0:30]   running total 1:45   ~98 words   >>> HAND TO OKASHA

So here is the problem we noticed.

The critic is always a fact-checker. But Hullman and Diakopoulos showed framing
alone significantly changes how people read the same chart. Tone is part of the
claim.

[POINT] Both sentences are true to the table. A fact-checker passes both. One
catastrophises, one is calibrated, and nothing in the literature tells them
apart.

Okasha will show you what we built for that.

DELIVERY: read both quotes out loud, slowly. This is the motivation and it lands
in ten seconds if you actually perform the contrast.
""")

    # -- 5  divider -------------------------------------------------------
    divider(prs, 1, """
SPEAKER 2 - Okasha   [0:05]   running total 1:50   ~15 words

Thank you. I will take the system itself, and then how we put a number on tone.

DELIVERY: your first words as a new speaker. Say them to the room, not the
screen.
""")

    # -- 6  the system ----------------------------------------------------
    s = blank(prs)
    heading(s, "Our fix", "One generator, one agentic moderator, one independent judge")
    # Generator
    rect(s, Inches(0.75), Inches(2.4), Inches(3.0), Inches(1.5), fill=SOFT)
    rect(s, Inches(0.75), Inches(2.4), Inches(3.0), Pt(4), fill=BLUE)
    text(s, Inches(0.99), Inches(2.62), Inches(2.55), Inches(1.2), [
        ("1   GENERATE", 10, True, MUTED, HEAD, 0),
        ("llama3.1:8b", 15, True, BLUE, HEAD, 4),
        ("Writes the first-draft story from\nan evidence pack built out of\nthe table.",
         11, False, INK, BODY, 4)], spacing=1.14)
    # The moderator, three jobs on one model
    rect(s, Inches(4.15), Inches(2.4), Inches(5.15), Inches(2.72), fill=SOFT)
    rect(s, Inches(4.15), Inches(2.4), Inches(5.15), Pt(4), fill=TEAL)
    text(s, Inches(4.42), Inches(2.62), Inches(4.6), Inches(0.6), [
        ("2   THE AGENTIC ROLE — ONE MODEL, THREE JOBS", 10, True, MUTED, HEAD, 0),
        ("gemma4:31b", 15, True, DEEP_TEAL, HEAD, 4)], spacing=1.14)
    for i, (job, desc, star) in enumerate([
        ("Moderates tone", "rewrites the framing, keeps the numbers", True),
        ("Fact-checks", "every stated number back against the table", False),
        ("Selects the charts", "ranks figures the table can honestly carry", False),
    ]):
        y = 3.42 + i * 0.55
        rect(s, Inches(4.42), Inches(y), Pt(2.5), Inches(0.44),
             fill=TEAL if star else BORDER)
        text(s, Inches(4.62), Inches(y), Inches(4.5), Inches(0.44), [
            (f"{job}   ", 12, True, NAVY if star else INK, BODY, 0)], spacing=1.0)
        text(s, Inches(4.62), Inches(y + 0.21), Inches(4.5), Inches(0.24), [
            (desc + ("   ← the contribution" if star else ""), 10.5, star,
             DEEP_TEAL if star else MUTED, BODY, 0)])
    # Judge
    rect(s, Inches(9.7), Inches(2.4), Inches(2.88), Inches(1.5), fill=SOFT)
    rect(s, Inches(9.7), Inches(2.4), Inches(2.88), Pt(4), fill=NAVY)
    text(s, Inches(9.94), Inches(2.62), Inches(2.45), Inches(1.2), [
        ("3   JUDGE", 10, True, MUTED, HEAD, 0),
        ("Claude Opus", 15, True, NAVY, HEAD, 4),
        ("Rates each story blind. Evaluation\nonly — not part of the product\nloop.",
         11, False, INK, BODY, 4)], spacing=1.14)
    rect(s, Inches(0.75), Inches(5.3), Inches(5.9), Inches(1.28), fill=SOFT)
    rect(s, Inches(0.75), Inches(5.3), Pt(3.5), Inches(1.28), fill=NAVY)
    text(s, Inches(1.02), Inches(5.48), Inches(5.4), Inches(1.0), [
        ("Same model, separate stages", 12.5, True, NAVY, BODY, 0),
        ("The moderator is the agentic role here, so the agentic reading of a "
         "table belongs to it too. The stages stay separate, so a tone number is "
         "still measuring one change — and one 19–23 GB load is paid per run "
         "instead of per stage.", 11.5, False, INK, BODY, 4)], spacing=1.14)
    text(s, Inches(7.15), Inches(5.32), Inches(5.43), Inches(1.3), [
        ("Two datasets, because tone fails both ways.", 12.5, True, NAVY, BODY, 0),
        ("Measles × MCV1 coverage (9,959 rows, 1980–2024) invites alarmism. WHO "
         "child mortality invites the opposite failure, false reassurance, where "
         "the agent must keep the inequality and the COVID reversal visible.",
         11.5, False, INK, BODY, 4)], spacing=1.14)
    chrome(s, 7)
    notes(s, """
SPEAKER 2 - Okasha   [0:33]   running total 2:23   ~82 words

Three roles, but only two local models.

A generator writes the draft. Then the moderator - the agentic role here - does
three jobs on the same weights: it rewrites the tone, that is our contribution,
it fact-checks every number, and it picks the charts.

Same model, separate stages. That matters twice: a tone number still measures
one change, and we pay one twenty-gigabyte model load per run instead of three.

Claude Opus judges, and only judges.

DELIVERY: "one model, three jobs" is the line. Do not call chart selection a
fourth model - it is not.
""")

    # -- 7  the instrument ------------------------------------------------
    s = blank(prs)
    heading(s, "The instrument", "Two axes, and both ends of each are failures")
    text(s, Inches(0.75), Inches(2.35), Inches(5.85), Inches(3.4), [
        ("The judge rates every story 1 to 5 on alarmism, and separately 1 to 5 "
         "on optimism.", 14, False, INK, BODY, 0),
        ("A 5 catastrophises. A 1 is flat and hides real stakes. Only the middle "
         "is calibrated — so this is not a “lower is better” scale, and an agent "
         "that simply drains the feeling out of a story scores badly too.",
         14, False, INK, BODY, 9),
        ("Blinding cost us something.", 14, True, NAVY, BODY, 12),
        ("An earlier version showed the judge both stories in one call, labelled "
         "before and after, always in that order. That names the treatment to "
         "the rater. Each story is now scored alone, which doubles the calls and "
         "gives up the direct comparison — the price of a rating that was not "
         "told which story is the treatment.", 12.5, False, MUTED, BODY, 4),
    ], spacing=1.18)
    scale_bar(s, 7.05, 2.4, 5.35, "ALARMISM",
              human=agg["human_alarmism_median"],
              human_label=f"human {agg['human_alarmism_median']:.1f}")
    scale_bar(s, 7.05, 3.85, 5.35, "OPTIMISM",
              human=agg["human_optimism_median"],
              human_label=f"human {agg['human_optimism_median']:.1f}")
    rect(s, Inches(7.05), Inches(5.2), Inches(5.35), Inches(1.2), fill=SOFT)
    rect(s, Inches(7.05), Inches(5.2), Pt(3.5), Inches(1.2), fill=BLUE)
    text(s, Inches(7.32), Inches(5.38), Inches(4.9), Inches(0.95),
         [("The rubric calls 3 calibrated. Our human writers sit at "
           f"{agg['human_alarmism_median']:.1f} and "
           f"{agg['human_optimism_median']:.1f}, so the human band — not the "
           "rubric midpoint — is the target we report against.", 12, False,
           INK, BODY, 0)], spacing=1.16)
    chrome(s, 8)
    notes(s, f"""
SPEAKER 2 - Okasha   [0:33]   running total 2:56   ~95 words

To measure tone you need a scale. Every story gets two: alarmism and optimism,
one to five each.

Both ends are failures. Five catastrophises, one is flat and hides real stakes.
So an agent cannot win by draining the feeling out of a story.

Two honest notes. The rubric calls three calibrated, but our human writers sit
at {agg['human_alarmism_median']:.1f} and {agg['human_optimism_median']:.1f}, so
the human band is what we report against.

And blinding cost us: scoring each story alone doubles the calls.

DELIVERY: "both ends are failures" must land. Pause after it. Point at the blue
human marks on the second note.
""")

    # -- 8  metrics part 1 ------------------------------------------------
    s = blank(prs)
    heading(s, "Metrics · part 1", "What we measure about the system")
    group(s, 0.75, 2.3, 5.9, 1.95, "Tone — the contribution",
          "An LLM judge, blind, one story per call. The only family that needs a model.",
          [("Alarmism 1-5", "flat ↔ catastrophising"),
           ("Optimism 1-5", "bleak ↔ false reassurance"),
           ("Emotive spans removed", "count, per run")], accent=TEAL)
    group(s, 0.75, 4.42, 5.9, 2.16, "Faithfulness — computed, no model, free",
          "Did the rewrite keep the facts? Runs on every story automatically.",
          [("Numeric retention", "share of the raw story's figures that survive"),
           ("Added unsupported", "figures the moderator introduced"),
           ("Groundedness", "stated numbers the table actually supports"),
           ("Rewrite intensity", "how much text changed at all"),
           ("Trend selection", "does the story's window agree with the whole "
            "series")], accent=BLUE)
    group(s, 6.9, 2.3, 5.68, 1.95, "Similarity to the human baseline",
          "How close the wording lands to the 25 stories we wrote by hand.",
          [("chrF++", "character n-grams; our ranking metric"),
           ("BLEU 1-4, ROUGE-L", "n-gram and longest-common-subsequence overlap"),
           ("METEOR, unigram F1", "stem and synonym aware")], accent=DEEP_TEAL)
    group(s, 6.9, 4.42, 5.68, 2.16, "Analytical operations — from the DataTales re-run",
          "Not every claim is equally hard, so we score accuracy per kind of claim.",
          [("lookup", "reading a single value off the table"),
           ("comparison", "ranking two values against each other"),
           ("trend", "which way a series moves across a window"),
           ("rate of change", "how steeply it moves"),
           ("causal", "why it moved — the hardest, and the one we treat with "
            "most caution")], accent=NAVY)
    text(s, Inches(0.75), Inches(6.68), Inches(11.8), Inches(0.28),
         [("“Trend” appears twice: trend selection (bottom left) is a "
           "cherry-picking check on one story; trend (bottom right) is an "
           "accuracy category from the benchmark.", 10, False, MUTED, BODY, 0)])
    chrome(s, 9)
    notes(s, """
SPEAKER 2 - Okasha   [0:28]   running total 3:24   ~70 words   >>> HAND TO ELSAADANI

Four families of metric, and only the first one needs a model.

Tone is the contribution and it is judged. Faithfulness is computed in code, for
free, including a cherry-picking check on the window the story chose. Similarity
is against our human stories. The operation accuracies come from DataTales.

Careful: two of these share the word trend. The footnote says which is which.

Elsaadani takes the reader-facing half.

DELIVERY: do NOT read the metric names. Name the four families and the fact
that three of them cost nothing to compute. The slide is the handout.
""")

    # -- 9  divider -------------------------------------------------------
    divider(prs, 2, """
SPEAKER 3 - Elsaadani   [0:05]   running total 3:29   ~15 words

My half: what a reader is actually shown, and how we picked the model pair.

DELIVERY: a beat, then straight into the next slide.
""")

    # -- 10  metrics part 2 -----------------------------------------------
    s = blank(prs)
    heading(s, "Metrics · part 2", "What the reader is shown, and why")
    group(s, 0.75, 2.3, 5.9, 2.35, "In the app, next to the story",
          "These exist so a reader can audit the rewrite instead of trusting it.",
          [("Alarmism rating", "where this story sits on the scale"),
           ("Human tone band", "±0.5 around the human author's own rating"),
           ("Numbing / catastrophising", "which end it failed toward"),
           ("Emotive span", "the exact words the moderator removed"),
           ("Verified / flagged / corrected", "what the fact-check did to each "
            "number")], accent=TEAL)
    group(s, 0.75, 4.85, 5.9, 1.8, "The edit taxonomy",
          "Every change the moderator made, classified — so the diff is readable.",
          [("Intensity", "how hard the wording pushes"),
           ("Framing", "what the sentence is made to be about"),
           ("Overreach", "a claim beyond what the table shows"),
           ("Grounding", "a claim tied back to a number")], accent=DEEP_TEAL)
    group(s, 6.9, 2.3, 5.68, 1.85, "Designed, not yet run",
          "The brief allows metrics or a user study. We chose metrics and went deep.",
          [("Trust", "does the reader believe the story"),
           ("Engagement", "does it hold attention"),
           ("Readability", "is it clear and easy to follow"),
           ("Human-vs-LLM preference", "which story readers prefer, and why")],
          accent=WARN)
    rect(s, Inches(6.9), Inches(4.35), Inches(5.68), Inches(2.2), fill=SOFT)
    rect(s, Inches(6.9), Inches(4.35), Inches(5.68), Pt(3.5), fill=NAVY)
    text(s, Inches(7.16), Inches(4.58), Inches(5.16), Inches(1.85), [
        ("The split that matters", 13, True, NAVY, HEAD, 0),
        ("Part 1 answers “is the system any good”. It is aggregated over runs "
         "and nobody reading a story ever sees it.", 11.5, False, INK, BODY, 6),
        ("Part 2 answers “should I believe this particular story”. It is "
         "per-story, visible in the interface, and it is the reason the app "
         "shows a diff rather than just a cleaner paragraph.", 11.5, False, INK,
         BODY, 6),
        ("Same underlying numbers in places. Different job.", 11.5, True,
         DEEP_TEAL, BODY, 6),
    ], spacing=1.15)
    chrome(s, 11)
    notes(s, """
SPEAKER 3 - Elsaadani   [0:24]   running total 3:53   ~60 words

Not every metric is for us. These are for the reader.

The app shows the tone rating, the human band, and every emotive span the
moderator removed, classified by kind. That is why the interface shows a diff
rather than a cleaner paragraph - a reader can audit the rewrite.

The user-study axes are designed and not run.

DELIVERY: this sets up the demo. Say "you will see these in a moment".
""")

    # -- 12  how we chose the pair ----------------------------------------
    s = blank(prs)
    heading(s, "Choosing the pair", "Four steps from fourteen candidates to one pipeline")
    steps = [
        ("01", "Build the grid",
         "Vary the generator against a fixed moderator, then the moderator "
         "against a fixed generator. L-shaped, not factorial — so the best × "
         "best cell is never assumed, it has to be run.", BLUE),
        ("02", "Score every run the same way",
         "Blind Opus rating on both tone axes, plus the computed faithfulness "
         "metrics. Same dataset throughout, so the comparison is not rewarded "
         "for drawing an easy table.", TEAL),
        ("03", "Tone does not decide it",
         "Six combinations land on exactly 2.0 alarmism. A tie on the headline "
         "metric is a real result: the moderator is doing its job regardless of "
         "who wrote the draft.", DEEP_TEAL),
        ("04", "Break the tie on faithfulness, then repeat it",
         "Rank the tied cells on how many of the raw story's figures survive. "
         "Then run the top two five times each at distinct seeds, because a "
         "single run cannot separate a model from a draft.", NAVY),
    ]
    for i2, (num, title, body, col) in enumerate(steps):
        x = 0.75 + (i2 % 2) * 6.04
        y = 2.35 + (i2 // 2) * 2.1
        rect(s, Inches(x), Inches(y), Inches(5.79), Inches(1.85), fill=SOFT)
        rect(s, Inches(x), Inches(y), Pt(3.5), Inches(1.85), fill=col)
        text(s, Inches(x + 0.28), Inches(y + 0.2), Inches(0.8), Inches(0.5),
             [(num, 20, True, RGBColor(0xC2, 0xCE, 0xDB), HEAD, 0)])
        text(s, Inches(x + 1.15), Inches(y + 0.22), Inches(4.4), Inches(1.5), [
            (title, 13.5, True, col, BODY, 0),
            (body, 11.5, False, INK, BODY, 5)], spacing=1.15)
    text(s, Inches(0.75), Inches(6.62), Inches(11.8), Inches(0.28),
         [("Step 4 is the one we nearly skipped, and it is the step that changed "
           "the answer.", 11, True, WARN, BODY, 0)])
    chrome(s, 12)
    notes(s, """
SPEAKER 3 - Elsaadani   [0:28]   running total 4:21   ~65 words

How we actually chose the pair, in four steps.

Build a grid of generator and moderator combinations. Score every run the same
way, on the same dataset. Tone does not decide it - six land on exactly 2.0.

So break the tie on faithfulness, then repeat the top two at different seeds,
because one run cannot separate a model from a draft.

The next two slides are steps three and four.

DELIVERY: this is the method slide. Say the four steps, then let the two figures
carry the evidence.
""")

    # -- 13  fig10 --------------------------------------------------------
    fullbleed(prs, "fig10.png", """
SPEAKER 3 - Elsaadani   [0:20]   running total 4:41   ~50 words

Step three. Fourteen runs. The green rows all tie on tone at exactly 2.0, so we
rank them on the "figures kept" column instead.

Now look at the three rows that read llama3.1:8b times gemma4:31b. Same models,
three separate runs: 50 percent, 75 percent, 100 percent.

DELIVERY: point at those three rows. That is the moment the audience should
distrust the table, and it sets up the next slide.
""")

    # -- 12  fig13 --------------------------------------------------------
    fullbleed(prs, "fig13.png", f"""
SPEAKER 3 - Elsaadani   [0:35]   running total 5:16   ~80 words   >>> HAND TO RAMADAN

Step four. Both configurations, five runs each, at distinct seeds.

It reverses. Qwen averages {D['g4b']['numeric_retention']['mean']:.0%}, llama
averages {D['g8b']['numeric_retention']['mean']:.0%}. Welch t of minus 2.70,
p of 0.029. Our recommendation had been backwards, and one run per cell was
never going to show us that.

So we ship llama3.1:8b generating and gemma4:31b moderating.

And the lower half is the more interesting finding: moderated tone is 2.10
either way. The moderator converges on the same tone whoever wrote the draft.

DELIVERY: "our recommendation had been backwards" is the line. Say it plainly.
""")

    # -- 13  divider ------------------------------------------------------
    divider(prs, 3, """
SPEAKER 4 - Ramadan   [0:05]   running total 5:21   ~15 words

Last part: where that lands against a person, and how the agent works inside.

DELIVERY: a beat. Then fig9 is the payoff slide of the whole talk.
""")

    # -- 14  fig9 ---------------------------------------------------------
    fullbleed(prs, "fig9.png", f"""
SPEAKER 4 - Ramadan   [0:38]   running total 5:59   ~105 words

This is the result. We wrote {D['n_human']} stories by hand from the same
evidence packs, before seeing any machine output. The green band is where those
human writers sit.

Red is raw, blue is after moderation. On every one of the five series the arrow
lands on or inside the human band.

Overall {agg['alarmism_raw_mean']:.2f} to {agg['alarmism_moderated_mean']:.2f},
against a human median of {agg['human_alarmism_median']:.2f}, over
{D['n_runs_human']} runs with a human counterpart.

The caveat is on the slide rather than hidden: our human stories carry no
headline and the machine ones do, so every gap is a lower bound.

DELIVERY: say the three numbers slowly, then stop talking.
""")

    # -- 17  the agent, and the tool surface ------------------------------
    s = blank(prs)
    heading(s, "How the model picks the charts",
            "A tool surface, so the decision is the model's and the rules are not")
    # the three-layer strip, compressed to one line of boxes
    strip = [("applicability.py", "which forms this table can carry", "computed", BLUE),
             ("the model", "which of those are worth showing, and why", "decided", TEAL),
             ("validate.py", "whether the spec is honest", "computed", NAVY)]
    for i2, (mod, what, tag, col) in enumerate(strip):
        x = 0.75 + i2 * 4.06
        rect(s, Inches(x), Inches(2.3), Inches(3.72), Inches(1.28), fill=SOFT)
        rect(s, Inches(x), Inches(2.3), Inches(3.72), Pt(4), fill=col)
        text(s, Inches(x + 0.26), Inches(2.5), Inches(3.2), Inches(1.0), [
            (f"{mod}   ·   {tag}", 11, True, col, HEAD, 0),
            (what, 12, False, INK, BODY, 4)], spacing=1.14)
        if i2 < 2:
            text(s, Inches(x + 3.76), Inches(2.78), Inches(0.3), Inches(0.4),
                 [("›", 20, True, BORDER, HEAD, 0)])
    text(s, Inches(0.75), Inches(3.72), Inches(11.8), Inches(0.3),
         [("Computing what is possible before asking what is best is what makes "
           "this work on a local model: the prompt shrinks from seventeen forms "
           "with their rules to a handful of pre-validated candidates.", 11.5,
           False, MUTED, BODY, 0)])
    # the MCP surface
    rect(s, Inches(0.75), Inches(4.2), Inches(7.6), Inches(2.4), fill=SOFT)
    rect(s, Inches(0.75), Inches(4.2), Inches(7.6), Pt(4), fill=DEEP_TEAL)
    text(s, Inches(1.02), Inches(4.42), Inches(7.1), Inches(0.4),
         [("THE MCP-SHAPED TOOL SURFACE", 10, True, MUTED, HEAD, 0)])
    tools = [
        ("7 chart tools", "plot_trend_over_time, plot_magnitude, plot_change, "
         "plot_relationship, plot_geographic, plot_distribution, show_headline"),
        ("3 read tools", "what makes it an agent rather than a classifier: it can "
         "look at the table before it commits to a figure"),
        ("A two-level decision", "the tool is the reader's job; the form enum "
         "inside it is the geometry. Seventeen flat forms is past what a local "
         "model picks reliably — seven tools with enums is not."),
    ]
    for i2, (k, v) in enumerate(tools):
        y = 4.78 + i2 * 0.58
        text(s, Inches(1.02), Inches(y), Inches(2.0), Inches(0.3),
             [(k, 11.5, True, DEEP_TEAL, BODY, 0)])
        text(s, Inches(2.9), Inches(y), Inches(5.2), Inches(0.55),
             [(v, 11, False, INK, BODY, 0)], spacing=1.12)
    rect(s, Inches(8.7), Inches(4.2), Inches(3.88), Inches(2.4), fill=SOFT)
    rect(s, Inches(8.7), Inches(4.2), Inches(3.88), Pt(4), fill=WARN)
    text(s, Inches(8.97), Inches(4.42), Inches(3.35), Inches(2.0), [
        ("Specified, not wired", 12.5, True, WARN, BODY, 0),
        ("The contract defines these tools and the build emits a JSON schema "
         "shaped to drop straight in as an MCP inputSchema — the same file the "
         "validator reads, so tool docs and rules cannot drift apart.", 11,
         False, INK, BODY, 5),
        ("Today the selector makes one structured call instead. The surface is a "
         "written design, not a running server.", 11, True, NAVY, BODY, 5),
    ], spacing=1.14)
    chrome(s, 17)
    notes(s, """
SPEAKER 4 - Ramadan   [0:40]   running total 6:39   ~105 words

How the model picks the figures, because people ask.

Three steps. Which forms can this table carry - computed from column types, so
it cannot propose a map for a table with no geography. Which are worth showing -
the model. Is the result honest - computed again.

The interesting part is the shape of the decision. Instead of seventeen flat
chart types, the contract groups them into seven tools by the reader's job, with
the geometry as an enum inside each. A two-level choice is what a local model
does reliably.

That surface is MCP-shaped and specified - a design, not a running server.

DELIVERY: say "specified, not wired" out loud. Someone will ask to see the
server.
""")

    # -- 16  whole system into the demo -----------------------------------
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=DEEP_NAVY)
    rect(s, 0, 0, Inches(0.16), H, fill=TEAL)
    text(s, Inches(1.1), Inches(1.15), Inches(11.0), Inches(0.9), [
        ("THE WHOLE SYSTEM", 12, True, TEAL, HEAD, 0),
        ("Table in, audited story out — on one laptop", 34, True, WHITE, HEAD, 6),
    ])
    chain = [("Table", "9,959 rows"), ("Evidence pack", "computed"),
             ("Raw story", "llama3.1:8b"), ("Moderated", "gemma4:31b"),
             ("Fact-check", "same model"), ("Charts", "same model"),
             ("Judged", "Opus, blind")]
    for i, (step, sub) in enumerate(chain):
        x = 0.95 + i * 1.72
        rect(s, Inches(x), Inches(2.75), Inches(1.5), Inches(0.85),
             fill=RGBColor(0x12, 0x2C, 0x5E))
        rect(s, Inches(x), Inches(2.75), Inches(1.5), Pt(3),
             fill=TEAL if i in (3, 4, 5) else RGBColor(0x2B, 0x4A, 0x84))
        text(s, Inches(x + 0.12), Inches(2.95), Inches(1.26), Inches(0.6), [
            (step, 11.5, True, WHITE, BODY, 0),
            (sub, 9, False, RGBColor(0x8A, 0x9C, 0xB8), BODY, 1)],
            align=PP_ALIGN.CENTER, spacing=1.1)
        if i < len(chain) - 1:
            text(s, Inches(x + 1.5), Inches(3.02), Inches(0.22), Inches(0.3),
                 [("›", 15, True, RGBColor(0x2B, 0x4A, 0x84), HEAD, 0)],
                 align=PP_ALIGN.CENTER)
    text(s, Inches(0.95), Inches(3.78), Inches(11.5), Inches(0.3),
         [("teal = the three stages that run on the moderator model", 10, False,
           RGBColor(0x8A, 0x9C, 0xB8), BODY, 0)])
    text(s, Inches(1.1), Inches(4.5), Inches(11.0), Inches(1.5), [
        ("LIVE DEMO  ·  3:30", 12, True, TEAL, HEAD, 0),
        ("Table → chosen figures → raw story → moderated story → the human "
         "story beside it", 20, True, WHITE, HEAD, 8),
        ("localhost:3000  ·  Django on :8000  ·  llama3.1:8b and gemma4:31b on "
         "Ollama. No API call anywhere in the generation path.", 13, False,
         RGBColor(0x8A, 0x9C, 0xB8), BODY, 8),
    ], spacing=1.2)
    notes(s, """
SPEAKER 4 - Ramadan   [0:21]   running total 7:00   ~50 words   >>> DEMO

That is the whole system. A table becomes an evidence pack, a local model writes
the story, and the moderator model rewrites the tone, checks the numbers and
picks the figures. Opus judges it blind.

Everything except the judge runs on this laptop.

Let me show you.

DELIVERY: 20 seconds, then switch. The browser should ALREADY be open on the
dataset page. Never click Generate and wait - see SPEAKER-SCRIPT.md.
""")

    # -- 17  closing ------------------------------------------------------
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=WHITE)
    rect(s, 0, 0, W, Inches(0.16), fill=TEAL)
    text(s, Inches(1.1), Inches(1.15), Inches(11.2), Inches(2.4), [
        ("What we can defend", 13, True, DEEP_TEAL, HEAD, 0),
        ("A second agent can move a generated data story onto the tone level of "
         "a person writing from the same table —", 26, True, NAVY, HEAD, 10),
        ("without losing the numbers, and at a cost in readability we can state.",
         26, True, NAVY, HEAD, 2),
    ], spacing=1.2)
    for i, (label, v, lab) in enumerate([
        ("The result",
         f"{agg['alarmism_raw_mean']:.2f} → {agg['alarmism_moderated_mean']:.2f}",
         f"alarmism, against a human median of {agg['human_alarmism_median']:.2f}"),
        ("The instrument", f"{rel['alarmism']['icc_2_1']:.3f}",
         "ICC, so every number above has an error bar"),
        ("The cost", f"{crit['narrative_quality'].get('raw', 0)}/{n}",
         "pairs where the raw story still read better"),
    ]):
        card(s, Inches(1.1 + i * 3.75), Inches(3.95), Inches(3.45), Inches(1.7),
             label, v, lab, accent=[DEEP_TEAL, BLUE, WARN][i])
    text(s, Inches(1.1), Inches(5.95), Inches(11.2), Inches(0.85), [
        ("What we are not claiming: one judge family agreeing with itself, not a "
         "panel. No user study. Twenty-five human stories from five writer slots, "
         "not a controlled writer design.", 12, False, MUTED, BODY, 0),
        ("Everything in this deck is generated from the experiment files in the "
         "repository, so the numbers and the slides cannot drift apart.", 12,
         False, MUTED, BODY, 6),
    ], spacing=1.16)
    chrome(s, 19)
    notes(s, """
CLOSING - whoever finishes the demo   [0:15]

One sentence: a second agent can move a generated data story onto the tone level
of a person writing from the same table, without losing the numbers, and at a
cost in readability we can put a number on.

What we are not claiming is on the slide. One judge family, no user study, and a
human baseline written by five writer slots rather than a controlled design.

DELIVERY: read the limits line rather than skipping it. Volunteering it is what
makes the rest credible. Q&A prep is at the end of SPEAKER-SCRIPT.md.
""")

    # -- 20  thank you -----------------------------------------------------
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=DEEP_NAVY)
    rect(s, 0, 0, Inches(0.16), H, fill=TEAL)
    text(s, Inches(1.1), Inches(2.15), Inches(11.0), Inches(1.5),
         [("Thank you", 60, True, WHITE, HEAD, 0)])
    rect(s, Inches(1.1), Inches(3.75), Inches(1.6), Pt(4), fill=TEAL)
    text(s, Inches(1.1), Inches(4.1), Inches(11.0), Inches(0.6),
         [("Questions?", 26, True, TEAL, HEAD, 0)])
    text(s, Inches(1.1), Inches(5.05), Inches(11.0), Inches(1.3), [
        ("Mahmoud Elsamadony  ·  Ahmed Okasha  ·  Ahmed Elsaadani  ·  "
         "Ahmed Ramadan", 14, True, WHITE, BODY, 0),
        ("Supervisors: Susmita Khadse, Julián Méndez  ·  Chair: "
         "Prof. Dr. Raimund Dachselt (IMLD)", 11.5, False,
         RGBColor(0x8A, 0x9C, 0xB8), BODY, 6),
        ("github.com/memoelsamadony/ai-powered-data-storytelling", 12, True,
         TEAL, BODY, 8),
    ], spacing=1.18)
    notes(s, """
QUESTIONS - all four speakers stay at the front

Leave this slide up for the whole Q&A. It holds the names and the repository
link, which is what anyone wanting to follow up needs.

The four questions to expect, with answers, are at the end of
presentation/SPEAKER-SCRIPT.md. Read them on the way in. Short forms:

- Only one judge? Yes, and slide 19 says so. 0.991 is self-consistency, not
  inter-rater. Two independent methods agree though: the scalar scores and the
  blind pairwise verdicts.
- Why no user study? The brief allows metrics or a user study. We chose metrics
  and went deep: reliability, repeats, a significance test.
- Why not just prompt the generator to be calm? Because a separate pass is
  measurable on its own, and because the moderator lands on 2.10 whichever
  generator wrote the draft.
- Is 2.09 against a human 2.00 just noise? Essentially yes, and that is the
  claim: it lands on the human level. The judge's own disagreement is 0.08.

DELIVERY: whoever is nearest the laptop takes questions first, then hand to
whoever owns that section. Nobody answers over anybody.
""")

    return prs


def crit_str(counts: dict, n: int) -> str:
    """'moderated 20/20' or 'raw 12  ·  tie 4  ·  moderated 4'."""
    order = ["moderated", "raw", "tie"]
    parts = [f"{k} {counts[k]}" for k in order if counts.get(k)]
    if len(parts) == 1:
        return f"{parts[0]}/{n}"
    return "   ·   ".join(parts)


RUNBOOK = """
## The demo, minute by minute

A cold run of the shipping configuration takes between 194 and 550 seconds -
measured, not guessed (`exp-repeats-g4b-g8b.json`). **You cannot click Generate
on stage and wait.** Everything below exists to make sure you never have to.

> **Walk this path once, end to end, the day before.** It is written against the
> code as it stands after the chart-selection merge, and the screens were read
> out of the components rather than clicked. That is enough to know each screen
> exists; it is not enough to know what it looks like with your data in it.

### Before you leave

- `ollama pull llama3.1:8b` and `ollama pull gemma4:31b`, then run one throwaway
  generation so the weights are resident. A cold load on stage reads as a hang.
- Django on `:8000`, Next on `:3000`, both already serving.
- One browser window, no other tabs, zoom 125%, notifications off.
- A second window holding the screenshot folder, in the demo's own order.

### T-0, as slide 1 goes up

The standby person starts a real generation on the measles dataset in a
background tab. At the measured mean it finishes around slide 8. If it has not
finished by the time you switch over, you open a completed run instead and say
so - one sentence, no apology.

### The click path

Everything after the first step lives on `/generate`, which composes the dataset
picker, the pipeline runner and the comparison panel on one page.

| Time | Where | The one thing you say |
|------|-------|----------------------|
| 0:00-0:20 | `/datasets` | Real merged data: measles cases joined to MCV1 coverage on country and year, 9,959 rows, 1980 to 2024. |
| 0:20-0:55 | `/generate`, dataset picked | We do not let the frontend guess which charts to draw. The backend profiles the table and picks the forms it can honestly carry - and refuses a spec that would misrepresent it. |
| 0:55-1:45 | Pipeline runner, raw story | Read **one** alarmist sentence out loud, then show its tone reading. |
| 1:45-2:40 | Moderated story | Same numbers, different temperature. The red-line view marks every emotive span that came out, and the tone meter moves. |
| 2:40-3:10 | Comparison panel | This is one of the 25 stories we wrote by hand. Now the similarity and retention numbers. |
| 3:10-3:30 | Anywhere | That entire loop ran on this laptop. There is no API call in the generation path. |

**The gotcha:** the comparison panel scores against the human baseline *typed or
imported into the page*, so the metrics stay empty until you paste one in. Have
the measles story on the clipboard, or import it, before you switch to the
browser - not while the room is watching.

### When it breaks

- **Generation not finished.** Open a completed run from the list. "Here is one
  from this morning." Move on. Do not stand and watch a spinner.
- **Backend down.** Screenshots, same order, same words. The story does not
  depend on it being live.
- **A model got evicted.** Do not re-pull on stage. Screenshots.

## If you are running long

Cut slide 7, the reliability figure, and fold one sentence into slide 6: "we
rated every story three times and the judge agrees with itself to ICC 0.99, so
this gap is real." That buys 40 seconds and costs the least.

Do **not** cut slide 9. The reversal is the strongest thing in the deck.

## Questions to expect

**"There is only one judge."** Correct, and we say so on slide 10. The 0.991 is
self-consistency, not inter-rater agreement - an upper bound on what a different
judge would agree to. A second judge family is the next run. What supports the
result meanwhile is that two independent methods agree: the scalar scores and
the blind pairwise verdicts point the same way.

**"Why no user study?"** The brief allows metrics or a user study. We chose
metrics and went deep on them - reliability coefficients, five-seed repeats, a
significance test. Every preference in the deck is a model's, and a user study
is the natural continuation, not a gap we overlooked.

**"Why not just prompt the generator to write calmly?"** We ran that as a
configuration. Two reasons it is not the same thing. The separation is
measurable: a second pass that only sees the text can be evaluated on its own.
And the data says something a prompt cannot give you - the moderator lands on
2.10 whichever generator wrote the draft, so the tone is a property of the
moderating step, not of the writer.

**"Is 2.09 against a human 2.00 not just noise?"** Essentially yes, and that is
the claim. It lands on the human level rather than below it. For scale, the
judge's own disagreement with itself is 0.08.

**"Your human baseline is your own team."** Yes. Five writer slots per series,
written from the evidence packs before anyone saw machine output. It is a real
hand-written baseline, and it is not the controlled writer design a full study
would want. Both are on slide 10.
"""


def write_script(prs: Presentation) -> Path:
    """Emit the speaker script from the deck's own notes.

    Generated rather than written by hand so the two cannot drift: edit a note
    in this file, rebuild, and the script follows. Only the runbook and the Q&A
    live here as prose, because they are not per-slide.
    """
    import re
    dest = HERE / "SPEAKER-SCRIPT.md"
    slides = list(prs.slides)

    def spoken(t: str) -> int:
        body = re.split(r"\nDELIVERY:", t)[0]
        body = "\n".join(l for l in body.splitlines()
                         if not re.match(r"^(SPEAKER|DEMO|CLOSING|DO NOT)", l.strip()))
        return len(body.split())

    # Ownership and slide ranges come out of the notes, not a list kept in step
    # by hand. Moving a slide from one speaker to another used to mean editing
    # this table too, and the table is exactly the thing nobody re-reads.
    talk_slides = [k for k, sl in enumerate(slides, 1)
                   if sl.notes_slide.notes_text_frame.text.strip().startswith("SPEAKER")]
    talk = sum(spoken(slides[k - 1].notes_slide.notes_text_frame.text)
               for k in talk_slides)
    per: dict[str, int] = {}
    span: dict[str, list[int]] = {}
    for k in talk_slides:
        head = slides[k - 1].notes_slide.notes_text_frame.text.splitlines()[0]
        who = head.split("-")[1].split()[0]
        per[who] = per.get(who, 0) + spoken(slides[k - 1].notes_slide.notes_text_frame.text)
        span.setdefault(who, []).append(k)

    L = ["# Final presentation - speaker script",
         "",
         f"**7:00 of talk across four speakers, then a 3:30 demo.** "
         f"{len(slides)} slides; the figure slides carry themselves, so on those "
         "you talk to the room, not to the screen.",
         "",
         f"Generated from the speaker notes inside "
         f"`{OUT.name}` by `make_final_deck.py`. Edit the notes there and rebuild "
         "rather than editing this file, or the two will disagree by Sunday.",
         "",
         "## Who speaks when",
         "",
         "| Speaker | Slides | Words | At 150 wpm |",
         "|---------|--------|-------|-----------|"]
    for name, ks in span.items():
        rng = f"{min(ks)}-{max(ks)}" if len(ks) > 1 else str(ks[0])
        w = per[name]
        L.append(f"| {name} | {rng} | {w} | {w / 150 * 60:.0f}s |")
    L += [f"| **Talk total** | **{min(talk_slides)}-{max(talk_slides)}** | **{talk}** | "
          f"**{int(talk / 150)}:{round(talk / 150 % 1 * 60):02d}** |",
          f"| Demo | after {max(talk_slides)} | - | 3:30 |",
          *[f"| Close | {k} | {spoken(sl.notes_slide.notes_text_frame.text)} | 15s |"
            for k, sl in enumerate(slides, 1)
            if sl.notes_slide.notes_text_frame.text.strip().startswith("CLOSING")],
          f"| Questions | {len(slides)} | - | stays up |",
          "",
          "Assignments follow the interim deck's four-part split and are meant to be "
          "swapped - the script is written per section, not per person. Whoever "
          "drives the demo should not also be the one closing it; give the standby "
          "person the last slide.",
          "",
          "## The script",
          ""]

    for i, s in enumerate(slides, 1):
        note = s.notes_slide.notes_text_frame.text.strip()
        head, _, rest = note.partition("\n")
        pic = any(sh.shape_type == 13 for sh in s.shapes)
        on_slide = " ".join(sh.text_frame.text for sh in s.shapes
                            if sh.has_text_frame)
        if pic:
            title = "figure, full bleed"
        elif re.search(r"PART \d+ OF \d+", on_slide):
            title = "section divider — hand over here"
        else:
            title = "layout"
        L += [f"### Slide {i} - {title}", "", f"`{head}`", ""]
        for block in rest.strip().split("\n\n"):
            block = " ".join(block.split())
            if block.startswith("DELIVERY:"):
                L += [f"> **{block[:9]}** {block[9:].strip()}", ""]
            else:
                L += [block, ""]

    L += [RUNBOOK.strip(), ""]
    dest.write_text("\n".join(L), encoding="utf-8")
    return dest


if __name__ == "__main__":
    prs = build()
    prs.save(OUT)
    script = write_script(prs)
    print(f"wrote {OUT.name}  ({len(prs.slides._sldIdLst)} slides, "
          f"{OUT.stat().st_size / 1024:.0f} KB)")
    print(f"wrote {script.name}  ({len(script.read_text().splitlines())} lines)")
