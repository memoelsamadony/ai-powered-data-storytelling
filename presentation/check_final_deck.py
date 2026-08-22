#!/usr/bin/env python3
"""Measure the built deck for text that runs past where it is allowed to.

A ``.pptx`` textbox does not clip. Give it more text than it has room for and
PowerPoint draws the overflow anyway, straight over the footer or the next
block, and the generator reports nothing wrong. There is no renderer here to
catch that by eye, so this measures it instead: lay each run out with the real
font file at the real point size, wrap it to the real box width, and add up the
line heights.

The fonts have to be the ones the deck actually asks for. Measuring Arial with
Helvetica's metrics would answer a question nobody asked.

    python3 presentation/check_final_deck.py

Exit code 1 if anything overruns, so it can gate a build.
"""

from __future__ import annotations

import sys
from pathlib import Path

from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).resolve().parent
DECK = HERE / "AI-Storytelling-Final-Talk.pptx"

EMU_IN = 914400.0
SLIDE_H = 7.5
FOOTER_TOP = 6.92          # the footer line sits here; nothing may reach it
SAFE_BOTTOM = 6.85

FONTS = {
    ("Arial", False): "/System/Library/Fonts/Supplemental/Arial.ttf",
    ("Arial", True): "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
    # Calibri is not installed system-wide on macOS; it ships inside the
    # PowerPoint bundle, which is where the real metrics are.
    ("Calibri", False): "/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts/Calibri.ttf",
    ("Calibri", True): "/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts/Calibrib.ttf",
}
_cache: dict = {}


def face(name: str, bold: bool, pt: float):
    """Pillow wants pixels; at 96 dpi a point is 96/72 px."""
    key = (name, bold, round(pt, 1))
    if key not in _cache:
        path = FONTS.get((name, bold))
        if path is None or not Path(path).exists():
            return None
        _cache[key] = ImageFont.truetype(path, int(round(pt * 96 / 72)))
    return _cache[key]


def wrapped_lines(txt: str, font, width_px: float) -> int:
    """Greedy wrap, the same rule PowerPoint uses for a word-wrapped box."""
    if not txt.strip():
        return 1
    lines, cur = 0, ""
    for word in txt.split():
        trial = f"{cur} {word}".strip()
        if font.getlength(trial) <= width_px or not cur:
            cur = trial
        else:
            lines += 1
            cur = word
    return lines + (1 if cur else 0)


def measure(shape) -> float | None:
    """Height in inches of the text as it will actually be drawn."""
    if not shape.has_text_frame:
        return None
    width_px = (shape.width / EMU_IN) * 96
    total = 0.0
    for para in shape.text_frame.paragraphs:
        runs = para.runs
        if not runs:
            continue
        pt = max((r.font.size.pt if r.font.size else 18) for r in runs)
        name = runs[0].font.name or "Calibri"
        bold = bool(runs[0].font.bold)
        f = face(name, bold, pt)
        txt = "".join(r.text for r in runs)
        if f is None:                       # unknown face: assume it is wider
            n_lines = max(1, int(len(txt) / max(1, width_px / (pt * 0.55)) + 0.999))
        else:
            n_lines = wrapped_lines(txt, f, width_px)
        spacing = para.line_spacing if isinstance(para.line_spacing, float) else 1.0
        before = para.space_before.pt if para.space_before else 0
        total += before / 72.0 + n_lines * (pt * 1.2 * spacing) / 72.0
    return total


def main() -> int:
    if not DECK.exists():
        print(f"no deck at {DECK}", file=sys.stderr)
        return 1
    missing = [k for k, v in FONTS.items() if not Path(v).exists()]
    if missing:
        print(f"WARNING: font files missing, results are estimates: {missing}")

    prs = Presentation(DECK)
    problems, checked = [], 0
    for i, slide in enumerate(prs.slides, 1):
        worst = 0.0
        for sh in slide.shapes:
            h = measure(sh)
            if h is None:
                continue
            checked += 1
            top = sh.top / EMU_IN
            bottom = top + h
            worst = max(worst, bottom)
            body = sh.text_frame.text.strip().replace("\n", " ")[:52]
            if bottom > SLIDE_H:
                problems.append((i, "OFF-SLIDE", bottom, body))
            elif bottom > FOOTER_TOP and top < FOOTER_TOP:
                problems.append((i, "HITS FOOTER", bottom, body))
            elif bottom > SAFE_BOTTOM and top < SAFE_BOTTOM:
                problems.append((i, "tight", bottom, body))
        n_pics = sum(1 for sh in slide.shapes if sh.shape_type == 13)
        kind = "figure" if n_pics else "layout"
        print(f"  slide {i:>2}  {kind:<7} lowest text ends at {worst:.2f}in "
              f"of {SLIDE_H}")

    print(f"\n{checked} text blocks measured across {len(prs.slides)} slides")
    if not problems:
        print("no overflow: every block clears the footer")
        return 0
    print(f"\n{len(problems)} problem(s):")
    for slide_no, kind, bottom, body in problems:
        print(f"  slide {slide_no:>2}  {kind:<11} ends {bottom:.2f}in  “{body}…”")
    return 1 if any(k != "tight" for _, k, _, _ in problems) else 0


if __name__ == "__main__":
    raise SystemExit(main())
